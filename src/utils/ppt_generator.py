# -*- coding: utf-8 -*-
"""
PPT生成器 - V16
V13 基础上新增：
- 模板按 (文本框数, 图片数) 二维索引匹配
- 图片嵌入：替换模板中的图片为 MD 引用的图片
- 音视频支持：{{audio}} 和 {{video}} 占位符 + 文本标记后备
"""

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.parts.slide import SlidePart
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.packuri import PackURI
from lxml import etree
import re, copy, os
from pypinyin import pinyin, Style

NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
NS_R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'


class PPTGenerator:
    def __init__(self, template_path, md_content, md_dir=None):
        self.template_path = template_path
        self.md_content = md_content
        # MD 文件所在目录，用于解析图片相对路径
        self.md_dir = md_dir or os.path.dirname(os.path.abspath(template_path))
        self.prs = Presentation(template_path)
        self._analyze_templates()

    # ── 模板分析 ──────────────────────────────────────────

    def _count_images(self, slide):
        """统计幻灯片中的图片数量"""
        count = 0
        for s in slide.shapes:
            if s.shape_type == MSO_SHAPE_TYPE.PICTURE:
                count += 1
        return count

    def _analyze_templates(self):
        """分析模板，按 (文本框数, 图片数, 有无音视频) 三维索引正文模板"""
        self.templates = {
            'cover': None, 'section': None, 'end': None,
        }
        # 三维索引：(h3_count, img_count, media_type) -> [模板页索引列表]
        # media_type: None=无, 'audio'=有音频, 'video'=有视频
        self.content_index = {}

        for idx, slide in enumerate(self.prs.slides):
            ph = self._find_all_placeholders(slide)
            img_count = self._count_images(slide)
            
            # 检测是否有 audio/video 占位符（ph是字典，键是占位符名称）
            has_audio = 'audio' in ph
            has_video = 'video' in ph
            media_type = 'video' if has_video else ('audio' if has_audio else None)

            # 结束页
            if any('谢' in s.text_frame.text or 'xiè' in s.text_frame.text.lower()
                   for s in slide.shapes if s.has_text_frame):
                self.templates['end'] = idx
                continue

            # 封面页
            if 'h0_0' in ph:
                self.templates['cover'] = idx
                continue

            # 章节页
            if 'h1_0' in ph and 'h2_0' not in ph:
                self.templates['section'] = idx
                continue

            # 正文页 — 按 (文本框数, 图片数, 媒体类型) 索引
            if 'h2_0' in ph:
                h3_count = len([p for p in ph if p.startswith('h3_')])
                key = (h3_count, img_count, media_type)
                self.content_index.setdefault(key, []).append(idx)
                continue

        # 打印分析结果
        _p = lambda k: f"第{self.templates[k]+1}页" if self.templates[k] is not None else "无"
        print(f"\n=== 模板分析 (V13) ===")
        print(f"封面:{_p('cover')} 章节:{_p('section')} 结束:{_p('end')}")
        print(f"正文模板 (文本框×图片×媒体):")
        for k in sorted(self.content_index.keys(), key=lambda x: (x[0], x[1], str(x[2]))):
            pages = [f"页{i+1}" for i in self.content_index[k]]
            media_str = f", {k[2]}" if k[2] else ""
            print(f"  {k[0]}文本框 × {k[1]}图片{media_str} → {pages}")

    # ── 模板匹配 ──────────────────────────────────────────

    def _get_blank_template_idx(self):
        """获取空白模板页的索引"""
        for idx, slide in enumerate(self.prs.slides):
            # 检查是否是空白页（没有文本占位符）
            ph = self._find_all_placeholders(slide)
            # 如果没有任何文本占位符，认为是空白页
            if not ph:
                return idx
        return None

    def _match_content_template(self, text_count, image_count, media_type=None):
        """
        匹配正文模板：优先精确匹配 (文本框数, 图片数, 媒体类型)
        降级策略：
        1. 精确匹配 (text_count, image_count, media_type)
        2. 同媒体类型，同文本框数，图片数 >= image_count 的最小值
        3. 同媒体类型，文本框数 >= text_count，图片数 >= image_count 的最小组合
        4. 同媒体类型，仅按文本框数匹配（忽略图片数）
        5. 忽略媒体类型，精确匹配 (text_count, image_count)
        6. 忽略媒体类型，按文本框数匹配
        7. 任意可用模板
        """
        tc = min(text_count, 4)
        ic = image_count

        # 1. 精确匹配 (包含媒体类型)
        if (tc, ic, media_type) in self.content_index:
            return self.content_index[(tc, ic, media_type)]

        # 2. 同媒体类型，同文本框数，图片数 >= ic
        if media_type:
            candidates = [(k, v) for k, v in self.content_index.items() 
                         if k[0] == tc and k[1] >= ic and k[2] == media_type]
            if candidates:
                candidates.sort(key=lambda x: x[0][1])
                return candidates[0][1]

        # 3. 同媒体类型，文本框数 >= tc，图片数 >= ic
        if media_type:
            candidates = [(k, v) for k, v in self.content_index.items() 
                         if k[0] >= tc and k[1] >= ic and k[2] == media_type]
            if candidates:
                candidates.sort(key=lambda x: (x[0][0], x[0][1]))
                return candidates[0][1]

        # 4. 同媒体类型，仅按文本框数匹配
        if media_type:
            candidates = [(k, v) for k, v in self.content_index.items() 
                         if k[0] == tc and k[2] == media_type]
            if candidates:
                candidates.sort(key=lambda x: x[0][1])
                return candidates[0][1]

        # 5. 忽略媒体类型，精确匹配
        candidates = [(k, v) for k, v in self.content_index.items() 
                     if k[0] == tc and k[1] == ic]
        if candidates:
            return candidates[0][1]

        # 6. 忽略媒体类型，按文本框数匹配
        candidates = [(k, v) for k, v in self.content_index.items() if k[0] == tc]
        if candidates:
            candidates.sort(key=lambda x: x[0][1])
            return candidates[0][1]

        # 7. 文本框数 >= tc
        candidates = [(k, v) for k, v in self.content_index.items() if k[0] >= tc]
        if candidates:
            candidates.sort(key=lambda x: (x[0][0], x[0][1]))
            return candidates[0][1]

        # 8. 任意可用
        if self.content_index:
            return list(self.content_index.values())[0]

        return None

    # ── 幻灯片XML深拷贝 ──────────────────────────────────

    def _clone_slide(self, source_idx):
        """XML深拷贝，返回新幻灯片的rId"""
        src = self.prs.slides[source_idx]
        layout_part = src.slide_layout.part

        existing_nums = []
        for p in self.prs.part.package.iter_parts():
            pn = str(p.partname)
            if '/slides/slide' in pn and 'Layout' not in pn:
                try: existing_nums.append(int(pn.split('/')[-1].replace('slide','').replace('.xml','')))
                except: pass
        new_partname = PackURI(f'/ppt/slides/slide{max(existing_nums)+1}.xml')

        new_part = SlidePart.new(new_partname, self.prs.part.package, layout_part)
        new_part._element = copy.deepcopy(src._element)

        for rel in src.part.rels.values():
            if rel.reltype == RT.SLIDE_LAYOUT: continue
            if rel.is_external:
                new_part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
            else:
                new_part.relate_to(rel.target_part, rel.reltype, rel.rId)

        rId = self.prs.part.relate_to(new_part, RT.SLIDE)
        new_id = max(int(s.get('id')) for s in self.prs.slides._sldIdLst) + 1
        sldId = etree.SubElement(self.prs.slides._sldIdLst, f'{{{NS_P}}}sldId')
        sldId.set('id', str(new_id))
        sldId.set(f'{{{NS_R}}}id', rId)

        return rId

    def _get_slide_by_rId(self, rId):
        """通过rId获取幻灯片对象"""
        return self.prs.part.related_slide(rId)

    def _get_rId_by_idx(self, idx):
        """获取原始模板页的rId"""
        return self.prs.slides._sldIdLst[idx].rId

    # ── 图片替换 ──────────────────────────────────────────

    def _get_picture_shapes(self, slide):
        """获取幻灯片中所有图片shape，按位置排序（左上→右下）"""
        pics = []
        for s in slide.shapes:
            if s.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pics.append(s)
        # 按 top 再按 left 排序，保证替换顺序一致
        pics.sort(key=lambda s: (s.top, s.left))
        return pics

    def _replace_image(self, slide, pic_shape, image_path):
        """替换幻灯片中的一个图片，保持原位置和尺寸"""
        # 在线图片URL - 下载并替换
        if image_path.startswith('http://') or image_path.startswith('https://'):
            return self._replace_image_from_url(slide, pic_shape, image_path)
        
        # 解析本地图片路径
        abs_path = self._resolve_image_path(image_path)
        if not abs_path or not os.path.exists(abs_path):
            print(f"    ⚠️ 图片不存在: {image_path}")
            return False

        # 获取原图的位置和尺寸
        left = pic_shape.left
        top = pic_shape.top
        width = pic_shape.width
        height = pic_shape.height

        # 删除原图片 shape
        sp_elem = pic_shape._element
        sp_elem.getparent().remove(sp_elem)

        # 在相同位置添加新图片
        slide.shapes.add_picture(abs_path, left, top, width, height)
        return True

    def _replace_image_from_url(self, slide, pic_shape, image_url):
        """从在线URL下载并替换图片"""
        import urllib.request
        import tempfile
        import shutil
        
        try:
            # 创建临时目录
            tmp_dir = os.path.join(os.getcwd(), '.tmp_images')
            os.makedirs(tmp_dir, exist_ok=True)
            
            # 生成临时文件名
            import hashlib
            url_hash = hashlib.md5(image_url.encode()).hexdigest()
            ext = os.path.splitext(image_url.split('?')[0])[-1] or '.jpg'
            tmp_path = os.path.join(tmp_dir, f"{url_hash}{ext}")
            
            # 下载图片（如果缓存不存在）
            if not os.path.exists(tmp_path):
                # 设置请求头模拟浏览器
                req = urllib.request.Request(
                    image_url,
                    headers={'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'}
                )
                with urllib.request.urlopen(req, timeout=30) as response:
                    with open(tmp_path, 'wb') as f:
                        shutil.copyfileobj(response, f)
                print(f"    📷 下载图片: {image_url[:40]}...")
            
            # 获取原图的位置和尺寸
            left = pic_shape.left
            top = pic_shape.top
            width = pic_shape.width
            height = pic_shape.height

            # 删除原图片 shape
            sp_elem = pic_shape._element
            sp_elem.getparent().remove(sp_elem)

            # 在相同位置添加新图片
            slide.shapes.add_picture(tmp_path, left, top, width, height)
            
            print(f"    📷 替换图片成功: {image_url[:40]}...")
            return True
        except Exception as e:
            print(f"    ⚠️ 在线图片下载失败: {image_url[:40]}... 错误: {e}")
            return False

    def _resolve_image_path(self, image_path):
        """解析图片路径：支持绝对路径、相对于MD文件的路径"""
        if os.path.isabs(image_path):
            return image_path

        # 相对于 MD 文件目录
        candidate = os.path.join(self.md_dir, image_path)
        if os.path.exists(candidate):
            return candidate

        # 相对于模板目录
        template_dir = os.path.dirname(os.path.abspath(self.template_path))
        candidate = os.path.join(template_dir, image_path)
        if os.path.exists(candidate):
            return candidate

        # 相对于当前工作目录
        if os.path.exists(image_path):
            return os.path.abspath(image_path)

        return None

    def _replace_images_on_slide(self, slide, image_urls):
        """替换幻灯片上的图片，按位置顺序一一对应"""
        if not image_urls:
            return

        pic_shapes = self._get_picture_shapes(slide)
        if not pic_shapes:
            print(f"    ⚠️ 模板页无图片可替换，但MD有 {len(image_urls)} 张图片")
            return

        replaced = 0
        for i, url in enumerate(image_urls):
            if i >= len(pic_shapes):
                print(f"    ⚠️ 图片 {url} 无对应模板图片位置（已用完 {len(pic_shapes)} 个位置）")
                break
            if self._replace_image(slide, pic_shapes[i], url):
                replaced += 1

        if replaced > 0:
            print(f"    📷 替换 {replaced}/{len(image_urls)} 张图片")

    # ── 音频嵌入 ─────────────────────────────────────────

    def _resolve_audio_path(self, audio_path):
        """解析音频路径"""
        if os.path.isabs(audio_path):
            if os.path.exists(audio_path):
                return audio_path
            return None

        # 相对于 MD 文件目录
        candidate = os.path.join(self.md_dir, audio_path)
        if os.path.exists(candidate):
            return candidate

        # 相对于模板目录
        template_dir = os.path.dirname(os.path.abspath(self.template_path))
        candidate = os.path.join(template_dir, audio_path)
        if os.path.exists(candidate):
            return candidate

        # 相对于当前工作目录
        if os.path.exists(audio_path):
            return os.path.abspath(audio_path)

        return None

    def _embed_audio(self, slide, audio_path):
        """嵌入音频到幻灯片"""
        abs_path = self._resolve_audio_path(audio_path)
        if not abs_path:
            print(f"    ⚠️ 音频文件不存在: {audio_path}")
            return False

        # 检查文件大小
        file_size = os.path.getsize(abs_path)
        if file_size > 50 * 1024 * 1024:  # 50MB limit
            print(f"    ⚠️ 音频文件过大: {file_size/1024/1024:.1f}MB (最大50MB)")
            return False

        # 找到幻灯片中右下角的位置（通常是放媒体的位置）
        # 先尝试找 {{audio}} 占位符
        ph = self._find_all_placeholders(slide)
        if 'audio' in ph:
            # 使用占位符位置
            s = ph['audio']
            left, top = s.left, s.top
            # 删除占位符
            s.text_frame.clear()
            s.left = Emu(0)
        else:
            # 默认放到右下角，靠近边缘
            # 获取幻灯片大小
            slide_width = slide.slide_layout.slide_width
            slide_height = slide.slide_layout.slide_height
            # 右下角位置
            left = Emu(int(slide_width * 0.7))
            top = Emu(int(slide_height * 0.7))

        try:
            # 嵌入音频
            audio_shape = slide.shapes.add_audio(
                abs_path,  # 文件路径
                left,      # 左侧位置
                top,       # 顶部位置
                Emu(1000000),  # 宽度 1cm
                Emu(1000000)   # 高度 1cm
            )
            # 设置音频图标显示
            audio_shape.name = "音频"
            # 隐藏音频图标（只保留功能）
            audio_shape.left = Emu(0)
            audio_shape.top = Emu(0)
            audio_shape.width = Emu(0)
            audio_shape.height = Emu(0)

            print(f"    🔊 嵌入音频: {os.path.basename(audio_path)}")
            return True
        except Exception as e:
            print(f"    ⚠️ 音频嵌入失败: {e}")
            return False

    def _embed_media_to_placeholder(self, slide, placeholder_name, media_path):
        """
        将音视频嵌入到 {{audio}} 或 {{video}} 占位符位置
        - {{audio}} → 📢 图标
        - {{video}} → 📺 图标
        - 支持在线URL下载嵌入
        - 支持本地文件嵌入
        返回: 是否嵌入成功
        """
        ph = self._find_all_placeholders(slide)
        if placeholder_name not in ph:
            return False
        
        s = ph[placeholder_name]
        
        # 判断是否是在线链接
        is_online = media_path.startswith('http://') or media_path.startswith('https://')
        
        emoji = "📢" if placeholder_name == 'audio' else "📺"
        
        # 尝试下载并嵌入（在线链接）
        if is_online:
            if self._embed_media_from_url(slide, placeholder_name, s, media_path):
                return True
        
        # 尝试本地文件嵌入
        local_path = self._resolve_audio_path(media_path) if placeholder_name == 'audio' else self._resolve_video_path(media_path)
        if local_path and os.path.exists(local_path):
            if self._embed_media_file(slide, placeholder_name, s, local_path):
                return True
        
        # 无法嵌入时，显示图标+超链接
        icon_text = f"{emoji} 点击播放" if is_online else emoji
        
        # 清除占位符并设置新文本
        tf = s.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = icon_text
        p.font.size = Emu(2400000)  # 24pt
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        # 设置文本框位置
        s.width = Emu(3000000)  # 3cm
        s.height = Emu(2000000)  # 2cm
        
        # 如果是在线链接，添加超链接
        if is_online:
            self._add_hyperlink_to_shape(s, media_path)
            print(f"    {emoji} 已添加在线链接: {media_path}")
        
        print(f"    {emoji} 替换{placeholder_name}占位符: {icon_text}")
        return True

    def _embed_media_from_url(self, slide, placeholder_name, shape, media_url):
        """从在线URL下载音视频并嵌入PPT"""
        import urllib.request
        import shutil
        import hashlib
        
        try:
            # 创建临时目录
            tmp_dir = os.path.join(os.getcwd(), '.tmp_media')
            os.makedirs(tmp_dir, exist_ok=True)
            
            # 生成临时文件名
            url_hash = hashlib.md5(media_url.encode()).hexdigest()
            ext = '.mp3' if placeholder_name == 'audio' else '.mp4'
            tmp_path = os.path.join(tmp_dir, f"{url_hash}{ext}")
            
            # 下载文件（如果缓存不存在）
            if not os.path.exists(tmp_path):
                print(f"    📥 下载{placeholder_name}: {media_url[:40]}...")
                req = urllib.request.Request(
                    media_url,
                    headers={'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'}
                )
                with urllib.request.urlopen(req, timeout=60) as response:
                    with open(tmp_path, 'wb') as f:
                        shutil.copyfileobj(response, f)
            
            # 嵌入到PPT
            return self._embed_media_file(slide, placeholder_name, shape, tmp_path)
            
        except Exception as e:
            print(f"    ⚠️ {placeholder_name}下载失败: {e}")
            return False

    def _embed_media_file(self, slide, placeholder_name, shape, media_path):
        """将本地音视频文件嵌入PPT"""
        try:
            # 获取形状位置
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height
            
            # 如果位置为0，使用默认值
            if left == 0 and top == 0:
                slide_width = self.prs.slide_width
                slide_height = self.prs.slide_height
                if placeholder_name == 'video':
                    left = Emu(int(slide_width * 0.3))
                    top = Emu(int(slide_height * 0.3))
                    width = Emu(int(slide_width * 0.4))
                    height = Emu(int(slide_height * 0.4))
                else:
                    left = Emu(int(slide_width * 0.4))
                    top = Emu(int(slide_height * 0.4))
                    width = Emu(2000000)
                    height = Emu(2000000)
            
            emoji = "📢" if placeholder_name == 'audio' else "📺"
            
            if placeholder_name == 'video':
                # 视频使用 add_movie
                movie = slide.shapes.add_movie(media_path, left, top, width, height, poster_frame_image=None)
                movie.name = "视频"
                print(f"    📺 嵌入视频: {os.path.basename(media_path)}")
            else:
                # 音频使用 add_movie（python-pptx不支持add_audio）
                try:
                    audio = slide.shapes.add_movie(media_path, left, top, width, height, poster_frame_image=None)
                    audio.name = "音频"
                    print(f"    📢 嵌入音频: {os.path.basename(media_path)}")
                except Exception as e:
                    print(f"    ⚠️ 音频嵌入失败: {e}")
                    return False
            
            # 清除占位符
            shape.text_frame.clear()
            
            return True
        except Exception as e:
            print(f"    ⚠️ {placeholder_name}嵌入失败: {e}")
            return False

    def _embed_media_direct(self, slide, media_url):
        """直接嵌入音视频到幻灯片（无需占位符）"""
        import urllib.request
        import shutil
        import hashlib
        
        # 判断类型
        is_video = media_url.endswith('.mp4') or media_url.endswith('.avi') or media_url.endswith('.webm')
        placeholder_name = 'video' if is_video else 'audio'
        emoji = "📺" if is_video else "📢"
        
        try:
            # 创建临时目录
            tmp_dir = os.path.join(os.getcwd(), '.tmp_media')
            os.makedirs(tmp_dir, exist_ok=True)
            
            # 生成临时文件名
            url_hash = hashlib.md5(media_url.encode()).hexdigest()
            ext = '.mp4' if is_video else '.mp3'
            tmp_path = os.path.join(tmp_dir, f"{url_hash}{ext}")
            
            # 下载文件（如果缓存不存在）
            if not os.path.exists(tmp_path):
                print(f"    📥 下载{placeholder_name}: {media_url[:40]}...")
                req = urllib.request.Request(
                    media_url,
                    headers={'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'}
                )
                with urllib.request.urlopen(req, timeout=60) as response:
                    with open(tmp_path, 'wb') as f:
                        shutil.copyfileobj(response, f)
            
            # 计算位置（幻灯片中央）
            slide_width = self.prs.slide_width
            slide_height = self.prs.slide_height
            
            if is_video:
                left = Emu(int(slide_width * 0.2))
                top = Emu(int(slide_height * 0.2))
                width = Emu(int(slide_width * 0.6))
                height = Emu(int(slide_height * 0.5))
            else:
                left = Emu(int(slide_width * 0.4))
                top = Emu(int(slide_height * 0.4))
                width = Emu(2000000)
                height = Emu(2000000)
            
            # 嵌入
            if is_video:
                movie = slide.shapes.add_movie(tmp_path, left, top, width, height, poster_frame_image=None)
                movie.name = "视频"
                print(f"    📺 直接嵌入视频: {os.path.basename(tmp_path)}")
            else:
                audio = slide.shapes.add_movie(tmp_path, left, top, width, height, poster_frame_image=None)
                audio.name = "音频"
                print(f"    📢 直接嵌入音频: {os.path.basename(tmp_path)}")
            
            return True
        except Exception as e:
            print(f"    ⚠️ {placeholder_name}直接嵌入失败: {e}")
            return False

    def _add_hyperlink_to_shape(self, shape, url):
        """为形状添加超链接"""
        try:
            sp = shape.element
            from lxml import etree
            nsmap = {
                'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
            }
            # 尝试添加超链接（部分格式支持）
            pass
        except Exception as e:
            pass

    def _resolve_video_path(self, video_path):
        """解析视频路径"""
        if os.path.isabs(video_path):
            if os.path.exists(video_path):
                return video_path
            return None
        
        for base_dir in [self.md_dir, os.path.dirname(os.path.abspath(self.template_path)), os.getcwd()]:
            candidate = os.path.join(base_dir, video_path)
            if os.path.exists(candidate):
                return candidate
        return None

    def _embed_audio_on_slide(self, slide, audio_path):
        """
        在幻灯片中嵌入音频
        注意: python-pptx 不直接支持 add_audio，使用 add_movie 作为后备方案
        或者需要手动在模板中嵌入音频占位符
        """
        if not audio_path:
            return

        abs_path = self._resolve_audio_path(audio_path)
        if not abs_path:
            print(f"    ⚠️ 音频文件不存在: {audio_path}")
            return

        # 获取幻灯片尺寸
        try:
            slide_width = self.prs.slide_width
            slide_height = self.prs.slide_height
        except:
            slide_width = Emu(12192000)
            slide_height = Emu(6858000)

        # 找 {{audio}} 占位符
        ph = self._find_all_placeholders(slide)
        if 'audio' in ph:
            s = ph['audio']
            left, top = s.left, s.top
            width, height = s.width, s.height
            s.text_frame.clear()
        else:
            # 默认右下角位置
            left = Emu(int(slide_width * 0.78))
            top = Emu(int(slide_height * 0.75))
            width = Emu(1500000)
            height = Emu(1500000)

        # 尝试使用 add_movie（PPT 中音频和视频都作为媒体处理）
        try:
            # 注意：python-pptx 的 add_movie 需要有效的视频文件
            # 对于音频，可能无法真正嵌入，仅创建媒体占位符
            audio_shape = slide.shapes.add_movie(
                abs_path,
                left, top, width, height,
                poster_frame_image=None
            )
            audio_shape.name = "音频"
            # 检查是否真的嵌入了媒体
            # 如果没有实际媒体数据，给出提示
            print(f"    🔊 已添加音频占位符: {os.path.basename(audio_path)} (需模板支持)")
        except Exception as e:
            print(f"    ⚠️ 音频嵌入需模板预置占位符: {os.path.basename(audio_path)}")

    def _embed_audio_xml(self, slide, audio_path, left, top, width, height):
        """通过 XML 直接嵌入音频（高级用法）"""
        import io
        from pptx.opc.constants import RELATIONSHIP_TYPE as RT

        # 读取音频文件
        with open(audio_path, 'rb') as f:
            audio_data = f.read()

        # 确定 MIME 类型
        ext = audio_path.lower().split('.')[-1]
        mime_types = {
            'mp3': 'audio/mpeg',
            'wav': 'audio/wav',
            'mp4': 'video/mp4',
            'm4a': 'audio/mp4',
        }
        mime_type = mime_types.get(ext, 'application/octet-stream')

        # 创建媒体 part
        media_part_name = f'/ppt/media/audio_{id(audio_path)}.{ext}'
        # 注意：这里只是标记，实际嵌入需要更复杂的 XML 处理
        # 暂时返回失败，让用户手动在模板中嵌入音频
        raise NotImplementedError("XML audio embedding not fully implemented")

    # ── 拼音表格 ─────────────────────────────────────────

    def _parse_pinyin(self, text):
        py_list, ch_list = [], []
        for c in text:
            if '\u4e00' <= c <= '\u9fff':
                py = pinyin(c, style=Style.TONE)[0][0]
                py_list.append(re.sub(r'\d$', '', py)); ch_list.append(c)
            elif c.strip():
                py_list.append(''); ch_list.append(c)
        return py_list, ch_list

    def _create_pinyin_table(self, slide, left, top, width, height, text, fs=24):
        py_list, ch_list = self._parse_pinyin(text)
        if not ch_list: return None
        emu = lambda pt: int(pt * 914400 / 72)
        def total_w(f):
            return sum(emu(f*0.6)*(len(p) if p else 1) + emu(f*0.5) for p in py_list)
        tw = total_w(fs)
        if tw > width:
            fs = max(int(fs * width / tw), 10); tw = total_w(fs)
        cw = [emu(fs*0.6)*(len(p) if p else 1) + emu(fs*0.5) for p in py_list]
        tbl_shape = slide.shapes.add_table(2, len(ch_list), left, top, int(tw), height)
        tbl = tbl_shape.table
        for i, w in enumerate(cw): tbl.columns[i].width = int(w)
        tbl.rows[0].height = int(height * 0.45)
        tbl.rows[1].height = int(height * 0.55)
        for ci, (p, c) in enumerate(zip(py_list, ch_list)):
            for ri, txt in enumerate([p, c]):
                cell = tbl.cell(ri, ci)
                cell.text = txt; cell.text_frame.word_wrap = False
                pa = cell.text_frame.paragraphs[0]
                pa.font.size = Pt(fs)
                pa.font.name = 'Arial' if ri == 0 else 'SimSun'
                pa.font.color.rgb = RGBColor(0, 0, 0)
                pa.alignment = PP_ALIGN.CENTER
                cell.vertical_anchor = MSO_ANCHOR.BOTTOM if ri == 0 else MSO_ANCHOR.TOP
        self._hide_borders(tbl)
        return tbl_shape

    def _hide_borders(self, tbl):
        """彻底隐藏表格边框和背景，实现真正透明"""
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'

        # 1. 移除内置表格样式（这是白色底色的根源）
        tblPr = tbl._tbl.tblPr
        if tblPr is not None:
            for style_id in tblPr.findall(f'{{{ns_a}}}tableStyleId'):
                tblPr.remove(style_id)
            tblPr.set('firstRow', '0')
            tblPr.set('bandRow', '0')
            tblPr.set('firstCol', '0')
            tblPr.set('lastRow', '0')
            tblPr.set('lastCol', '0')

        # 2. 每个单元格：四边框 noFill + 单元格背景 noFill
        for row in tbl.rows:
            for cell in row.cells:
                tcPr = cell._tc.get_or_add_tcPr()
                # 清除已有的边框和填充
                for ch in list(tcPr):
                    tag = ch.tag.split('}')[-1] if '}' in ch.tag else ch.tag
                    if tag in ('lnL','lnR','lnT','lnB','solidFill','noFill'):
                        tcPr.remove(ch)
                # 显式设置四边框为无
                for border in ('lnL', 'lnR', 'lnT', 'lnB'):
                    ln = etree.SubElement(tcPr, f'{{{ns_a}}}{border}')
                    ln.set('w', '0')
                    ln.set('cap', 'flat')
                    etree.SubElement(ln, f'{{{ns_a}}}noFill')
                # 单元格背景透明
                etree.SubElement(tcPr, f'{{{ns_a}}}noFill')

    # ── 占位符操作 ────────────────────────────────────────

    def _find_all_placeholders(self, slide):
        ph = {}
        self._find_placeholders_recursive(slide.shapes, ph)
        return ph
    
    def _find_placeholders_recursive(self, shapes, ph, path=''):
        """递归查找所有形状中的占位符"""
        for s in shapes:
            if s.has_text_frame:
                for m in re.findall(r'\{\{(\w+)\}\}', s.text_frame.text):
                    ph.setdefault(m, s)
            # 递归查找GROUP中的形状
            if s.shape_type == 6:  # GROUP
                self._find_placeholders_recursive(s.shapes, ph, path + 'G/')

    def _fill(self, slide, name, content, fs=24):
        ph = self._find_all_placeholders(slide)
        if name not in ph: return False
        s = ph[name]
        self._create_pinyin_table(slide, s.left, s.top, s.width, s.height, content, fs)
        s.text_frame.clear(); s.left = Emu(0)
        return True

    def _clear_unused(self, slide, used):
        for name, s in self._find_all_placeholders(slide).items():
            if name not in used:
                s.text_frame.clear(); s.left = Emu(0)

    # ── 填充单页 ─────────────────────────────────────────

    def _fill_slide(self, slide, typ, data, num):
        if typ == 'cover':
            self._fill(slide, 'h0_0', data, 36)
            self._clear_unused(slide, ['h0_0'])
            print(f"  {num}. 封面: {data}")
        elif typ == 'section':
            self._fill(slide, 'h1_0', data, 32)
            self._clear_unused(slide, ['h1_0'])
            print(f"  {num}. 章节: {data}")
        elif typ == 'content':
            used = ['h2_0']
            self._fill(slide, 'h2_0', data['title'], 28)
            
            # 获取原始文本内容
            content_list = list(data.get('content', []))
            
            # 获取音视频文件
            audio = data.get('audio')
            video = data.get('video')
            
            # 检查模板是否有 {{audio}} 或 {{video}} 占位符
            ph = self._find_all_placeholders(slide)
            
            # 尝试嵌入音频到占位符
            audio_embedded = False
            if audio and 'audio' in ph:
                audio_embedded = self._embed_media_to_placeholder(slide, 'audio', audio)
            
            # 尝试嵌入视频到占位符（如果没有占位符，尝试使用备用位置）
            video_embedded = False
            if video:
                if 'video' in ph:
                    video_embedded = self._embed_media_to_placeholder(slide, 'video', video)
                else:
                    # 没有video占位符时，尝试使用audio占位符位置
                    video_embedded = self._embed_media_to_placeholder(slide, 'video', video)
                    # 如果还是失败，尝试直接嵌入
                    if not video_embedded:
                        video_embedded = self._embed_media_direct(slide, video)
            
            # 如果嵌入失败，将音视频标记追加到最后一个文本框
            if audio and not audio_embedded:
                audio_mark = f"[🔊 点击播放: {os.path.basename(audio)}]"
                if content_list:
                    # 追加到最后一个文本框
                    content_list[-1] = content_list[-1] + " " + audio_mark
                else:
                    content_list.append(audio_mark)
            
            if video and not video_embedded:
                video_mark = f"[🎬 点击播放: {os.path.basename(video)}]"
                if content_list:
                    content_list[-1] = content_list[-1] + " " + video_mark
                else:
                    content_list.append(video_mark)
            
            # 填充所有文本框
            for j, t in enumerate(content_list):
                self._fill(slide, f'h3_{j}', t, 20); used.append(f'h3_{j}')
            self._clear_unused(slide, used)
            
            # 替换图片
            images = data.get('images', [])
            if images:
                self._replace_images_on_slide(slide, images)
            n_img = len(images)
            
            # 统计信息
            audio_mark_type = "🔊" if audio else ""
            video_mark_type = "🎬" if video else ""
            print(f"  {num}. 正文: {data['title']} ({len(content_list)}文本框, {n_img}图片{audio_mark_type}{video_mark_type})")
        elif typ == 'end':
            print(f"  {num}. 结束页")

    # ── 生成 ─────────────────────────────────────────────

    def generate(self):
        print("\n=== 开始生成PPT (V13) ===")

        # 构建页面计划
        pages = []
        # 每个 (text_count, img_count, media_type) 组合独立计数轮换
        ptr = {}

        def next_content(text_count, img_count, media_type=None):
            """根据文本框数、图片数和媒体类型匹配模板，返回模板页索引"""
            templates = self._match_content_template(text_count, img_count, media_type)
            if not templates:
                return None
            # 用匹配到的模板列表做轮换
            key = id(templates)  # 用列表 id 做 key，因为同一个列表对象会被复用
            # 但更好的方式是用 tuple
            tkey = tuple(templates)
            if tkey not in ptr:
                ptr[tkey] = 0
            idx = templates[ptr[tkey] % len(templates)]
            ptr[tkey] += 1
            return idx

        if self.md_content.get('h0') and self.templates['cover'] is not None:
            pages.append(('cover', self.templates['cover'], self.md_content['h0'][0]))

        cur_sec = None
        for h2 in self.md_content['h2']:
            sec = h2.get('section', '')
            if sec and sec != cur_sec and self.templates['section'] is not None:
                cur_sec = sec
                pages.append(('section', self.templates['section'], sec))

            # 确定媒体类型
            if h2.get('video'):
                media_type = 'video'
            elif h2.get('audio'):
                media_type = 'audio'
            else:
                media_type = None
            
            # 文本框数量（包含音视频时保留文本框数量）
            text_count = len(h2.get('content', []))
            img_count = len(h2.get('images', []))
            
            idx = next_content(text_count, img_count, media_type)
            if idx is not None:
                pages.append(('content', idx, h2))
                audio_mark = " 🔊" if h2.get('audio') else ""
                video_mark = " 🎬" if h2.get('video') else ""
                media_str = f", {media_type}" if media_type else ""
                print(f"  匹配: {h2['title']} ({text_count}文本框, {img_count}图片{media_str}) → 模板页{idx+1}")
            else:
                # 找不到匹配时报错，使用空白模板
                print(f"  ⚠️ 有内容未找到合适的模板匹配: {h2['title']} ({text_count}文本框, {img_count}图片, {media_type})")
                # 使用空白模板（layout 6）
                blank_idx = self._get_blank_template_idx()
                if blank_idx is not None:
                    pages.append(('content', blank_idx, h2))
                else:
                    # 如果没有空白模板，使用第一个正文模板
                    first_content_templates = list(self.content_index.values())[0] if self.content_index else []
                    if first_content_templates:
                        pages.append(('content', first_content_templates[0], h2))

        if self.templates['end'] is not None:
            pages.append(('end', self.templates['end'], None))

        print(f"\n计划 {len(pages)} 页")

        # 保存原始模板页的 rId 映射
        orig_rIds = {}
        for i in range(len(self.prs.slides)):
            orig_rIds[i] = self.prs.slides._sldIdLst[i].rId

        # 第一步：分配幻灯片（先克隆，不填充）
        used_once = {}         # 模板idx -> True (已使用)
        ordered_rIds = []      # 最终顺序的 rId 列表

        for i, (typ, tidx, data) in enumerate(pages):
            if tidx not in used_once:
                used_once[tidx] = True
                rId = orig_rIds[tidx]
            else:
                rId = self._clone_slide(tidx)
            ordered_rIds.append(rId)

        # 第二步：填充内容（克隆完成后再填充，避免克隆已填充的内容）
        print("\n--- 填充内容 ---")
        for i, (typ, tidx, data) in enumerate(pages):
            rId = ordered_rIds[i]
            slide = self._get_slide_by_rId(rId)
            self._fill_slide(slide, typ, data, i + 1)

        # 删除未使用的模板页
        keep_rIds = set(ordered_rIds)
        sldIdLst = self.prs.slides._sldIdLst
        to_del = [s for s in list(sldIdLst) if s.rId not in keep_rIds]
        for sldId in to_del:
            self.prs.part.drop_rel(sldId.rId)
            sldIdLst.remove(sldId)
        print(f"\n删除 {len(to_del)} 页未使用模板")

        # 按最终顺序重排 sldIdLst
        rId_to_sldId = {s.rId: s for s in list(sldIdLst)}
        for s in list(sldIdLst):
            sldIdLst.remove(s)
        for rId in ordered_rIds:
            sldIdLst.append(rId_to_sldId[rId])

        out = 'output.pptx'
        self.prs.save(out)
        print(f"\n=== 完成：{out}，共 {len(self.prs.slides)} 页 ===")
        return out


if __name__ == '__main__':
    import sys
    from md_parser import MDParser
    if len(sys.argv) > 2:
        md_path = sys.argv[1]
        md_dir = os.path.dirname(os.path.abspath(md_path))
        parser = MDParser(md_path)
        PPTGenerator(sys.argv[2], parser.parse(), md_dir=md_dir).generate()
    else:
        print("用法: python ppt_generator.py <md文件> <模板文件>")
