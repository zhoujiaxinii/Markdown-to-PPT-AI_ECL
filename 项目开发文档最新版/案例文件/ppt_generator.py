# -*- coding: utf-8 -*-
"""
PPT生成器 - 智能模板匹配版本 V2
正确处理：封面 → 目录 → 章节页 → 正文页
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
import re
from pypinyin import pinyin, Style

class PPTGenerator:
    def __init__(self, template_path, md_content):
        self.template_path = template_path
        self.md_content = md_content
        self.presentation = Presentation(template_path)
        self._analyze_templates()
    
    def _analyze_templates(self):
        """分析模板中每一页的特征"""
        self.templates = {
            'cover': None,      # 封面页: {{h0_0}}
            'toc': None,        # 目录页: 无占位符但有固定文字"目录"
            'section': None,    # 章节页: {{h1_0}}
            'end': None,        # 结束页: 含有"谢谢"
            'content': []       # 正文页: {{h2_0}} + {{h3_*}}
        }
        
        for idx, slide in enumerate(self.presentation.slides):
            placeholders = self._find_all_placeholders(slide)
            
            # 检查是否是结束页（含有"谢谢"）
            is_end_page = False
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text
                    if '谢' in text or 'xiè' in text.lower():
                        self.templates['end'] = {'idx': idx, 'slide': slide}
                        is_end_page = True
                        break
            
            if is_end_page:
                continue
            
            if 'h0_0' in placeholders:
                self.templates['cover'] = {'idx': idx, 'slide': slide}
            elif 'h1_0' in placeholders and 'h2_0' not in placeholders:
                self.templates['section'] = {'idx': idx, 'slide': slide}
            elif 'h2_0' in placeholders:
                text_count = len([p for p in placeholders if p.startswith('h3_')])
                self.templates['content'].append({
                    'idx': idx,
                    'text_count': text_count,
                    'slide': slide
                })
            else:
                # 检查是否有"目录"字样
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text_frame.text
                        if '目录' in text or 'mù lù' in text.lower():
                            self.templates['toc'] = {'idx': idx, 'slide': slide}
                            break
        
        # 按文本框数排序
        self.templates['content'].sort(key=lambda x: x['text_count'])
        
        print(f"\n=== 模板分析 ===")
        if self.templates['cover']:
            print(f"封面页: 第{self.templates['cover']['idx']+1}页")
        if self.templates['toc']:
            print(f"目录页: 第{self.templates['toc']['idx']+1}页")
        if self.templates['section']:
            print(f"章节页: 第{self.templates['section']['idx']+1}页")
        if self.templates['end']:
            print(f"结束页: 第{self.templates['end']['idx']+1}页")
        print(f"正文模板: {len(self.templates['content'])}个")
    
    def _find_all_placeholders(self, slide):
        """查找页面所有占位符"""
        placeholders = {}
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                matches = re.findall(r'\{\{(\w+)\}\}', text)
                for match in matches:
                    if match not in placeholders:
                        placeholders[match] = shape
        return placeholders
    
    def _find_best_content_template(self, text_count):
        """找到最佳匹配的正文模板"""
        candidates = self.templates['content']
        
        if not candidates:
            return None
        
        # 找文本框数刚好够用的最小模板
        for template in candidates:
            if template['text_count'] >= text_count:
                return template
        
        # 如果都不够，返回文本框最多的
        return max(candidates, key=lambda x: x['text_count'])
    
    def _parse_pinyin_chars(self, text):
        """解析文本，返回拼音和字符列表"""
        pinyin_list = []
        char_list = []
        
        for char in text:
            if '\u4e00' <= char <= '\u9fff':
                py = pinyin(char, style=Style.TONE)[0][0]
                py_clean = re.sub(r'(\d)$', '', py)
                pinyin_list.append(py_clean)
                char_list.append(char)
            elif char.strip():
                pinyin_list.append('')
                char_list.append(char)
        
        return pinyin_list, char_list
    
    def _create_pinyin_table(self, slide, left, top, width, height, text, font_size=24):
        """创建拼音-汉字对齐表格"""
        pinyin_list, char_list = self._parse_pinyin_chars(text)
        
        if not char_list:
            return None
        
        pt_to_emu = lambda pt: int(pt * 914400 / 72)
        
        def calc_total_width(fs):
            letter_width = pt_to_emu(fs * 0.5)
            total = 0
            for py in pinyin_list:
                py_len = len(py) if py else 1
                total += pt_to_emu(fs * 0.6) * py_len + letter_width
            return total
        
        total_width = calc_total_width(font_size)
        actual_font_size = font_size
        
        if total_width > width:
            scale = width / total_width
            actual_font_size = max(int(font_size * scale), 10)
            total_width = calc_total_width(actual_font_size)
        
        letter_width_emu = pt_to_emu(actual_font_size * 0.5)
        col_widths = []
        for py in pinyin_list:
            py_len = len(py) if py else 1
            col_widths.append(pt_to_emu(actual_font_size * 0.6) * py_len + letter_width_emu)
        
        table_shape = slide.shapes.add_table(2, len(char_list), left, top, int(total_width), height)
        table = table_shape.table
        
        for col_idx, col_width in enumerate(col_widths):
            table.columns[col_idx].width = int(col_width)
        
        table.rows[0].height = int(height * 0.45)
        table.rows[1].height = int(height * 0.55)
        
        for col_idx, (py, char) in enumerate(zip(pinyin_list, char_list)):
            pinyin_cell = table.cell(0, col_idx)
            pinyin_cell.text = py
            pinyin_cell.text_frame.word_wrap = False
            pinyin_para = pinyin_cell.text_frame.paragraphs[0]
            pinyin_para.font.size = Pt(actual_font_size)
            pinyin_para.font.name = 'Arial'
            pinyin_para.font.color.rgb = RGBColor(0, 0, 0)
            pinyin_para.alignment = PP_ALIGN.CENTER
            pinyin_cell.vertical_anchor = MSO_ANCHOR.BOTTOM
            
            char_cell = table.cell(1, col_idx)
            char_cell.text = char
            char_cell.text_frame.word_wrap = False
            char_para = char_cell.text_frame.paragraphs[0]
            char_para.font.size = Pt(actual_font_size)
            char_para.font.name = 'SimSun'
            char_para.font.color.rgb = RGBColor(0, 0, 0)
            char_para.alignment = PP_ALIGN.CENTER
            char_cell.vertical_anchor = MSO_ANCHOR.TOP
        
        self._hide_table_style(table)
        return table_shape
    
    def _hide_table_style(self, table):
        """隐藏表格边框和底色"""
        from lxml import etree
        
        tbl = table._tbl
        tblPr = tbl.tblPr
        if tblPr is None:
            tblPr = etree.SubElement(tbl, '{http://schemas.openxmlformats.org/drawingml/2006/main}tblPr')
        
        tblBorders = etree.SubElement(tblPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}tblBorders')
        for border_name in ['left', 'right', 'top', 'bottom', 'insideH', 'insideV']:
            border = etree.SubElement(tblBorders, '{http://schemas.openxmlformats.org/drawingml/2006/main}' + border_name)
            border.set('w', '0')
            etree.SubElement(border, '{http://schemas.openxmlformats.org/drawingml/2006/main}noFill')
        
        for row in table.rows:
            for cell in row.cells:
                tc = cell._tc
                tcPr = tc.get_or_add_tcPr()
                for child in list(tcPr):
                    if 'ln' in child.tag:
                        tcPr.remove(child)
                for border_name in ['lnL', 'lnR', 'lnT', 'lnB']:
                    ln = etree.SubElement(tcPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}' + border_name)
                    ln.set('w', '0')
                    etree.SubElement(ln, '{http://schemas.openxmlformats.org/drawingml/2006/main}noFill')
                solidFill = tcPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
                if solidFill is not None:
                    tcPr.remove(solidFill)
                etree.SubElement(tcPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}noFill')
    
    def _replace_placeholder_with_pinyin(self, slide, placeholder_type, content, font_size=24):
        """替换占位符为拼音表格"""
        for shape in slide.shapes:
            if shape.has_text_frame:
                if placeholder_type in shape.text_frame.text:
                    left, top, width, height = shape.left, shape.top, shape.width, shape.height
                    table_shape = self._create_pinyin_table(slide, left, top, width, height, content, font_size)
                    if table_shape:
                        shape.text_frame.clear()
                        shape.left = Emu(0)
                        shape.top = Emu(0)
                        shape.width = Emu(0)
                        shape.height = Emu(0)
                    return True
        return False
    
    def _replace_placeholder_text(self, slide, placeholder_type, content):
        """简单替换占位符文本"""
        for shape in slide.shapes:
            if shape.has_text_frame:
                if placeholder_type in shape.text_frame.text:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if placeholder_type in run.text:
                                run.text = run.text.replace('{{' + placeholder_type + '}}', content)
                    return True
        return False
    
    def _clear_placeholder(self, slide, placeholder_type):
        """清除指定占位符"""
        for shape in slide.shapes:
            if shape.has_text_frame:
                if placeholder_type in shape.text_frame.text:
                    shape.text_frame.clear()
                    shape.left = Emu(0)
                    shape.top = Emu(0)
                    shape.width = Emu(0)
                    shape.height = Emu(0)
                    return True
        return False
    
    def _clear_unused_placeholders(self, slide, used_list):
        """清除未使用的占位符"""
        placeholders = self._find_all_placeholders(slide)
        cleared = []
        for name in placeholders:
            if name not in used_list:
                self._clear_placeholder(slide, name)
                cleared.append(name)
        if cleared:
            print(f"    清除: {cleared}")
    
    def _duplicate_slide(self, slide_idx):
        """复制幻灯片"""
        source = self.presentation.slides[slide_idx]
        slide_layout = source.slide_layout
        new_slide = self.presentation.slides.add_slide(slide_layout)
        new_idx = len(self.presentation.slides) - 1
        return new_idx, new_slide
    
    def generate(self):
        """生成PPT"""
        print("\n=== 开始生成PPT ===")
        
        # 收集所有需要生成的页面
        pages_to_generate = []
        
        # 1. 封面页
        if self.md_content['h0']:
            pages_to_generate.append({
                'type': 'cover',
                'title': self.md_content['h0'][0]
            })
        
        # 2. 目录页（列出所有章节）
        if self.md_content['h1']:
            pages_to_generate.append({
                'type': 'toc',
                'sections': self.md_content['h1']
            })
        
        # 3. 章节页 + 正文页
        current_section = None
        section_index = 0
        
        for h2 in self.md_content['h2']:
            section = h2.get('section', '')
            
            # 如果是新章节，先添加章节页
            if section and section != current_section:
                current_section = section
                pages_to_generate.append({
                    'type': 'section',
                    'title': section
                })
                section_index += 1
            
            # 添加正文页
            pages_to_generate.append({
                'type': 'content',
                'title': h2['title'],
                'contents': h2.get('content', [])
            })
        
        print(f"需要生成 {len(pages_to_generate)} 页")
        for i, page in enumerate(pages_to_generate):
            print(f"  {i+1}. {page['type']}: {page.get('title', '')[:20]}")
        
        # 记录需要保留的幻灯片
        slides_to_keep = []
        
        # 生成每一页
        for i, page in enumerate(pages_to_generate):
            print(f"\n第{i+1}页: {page['type']} - {page.get('title', '')[:15]}...")
            
            if page['type'] == 'cover':
                # 封面页
                if self.templates['cover']:
                    template_idx = self.templates['cover']['idx']
                    slide = self.presentation.slides[template_idx]
                    self._replace_placeholder_with_pinyin(slide, 'h0_0', page['title'], font_size=36)
                    self._clear_unused_placeholders(slide, ['h0_0'])
                    slides_to_keep.append(template_idx)
                    print(f"  封面已填充: {page['title']}")
            
            elif page['type'] == 'toc':
                # 目录页
                if self.templates['toc']:
                    template_idx = self.templates['toc']['idx']
                    slide = self.presentation.slides[template_idx]
                    
                    # 填充目录项（假设目录页有多个文本框）
                    for j, section in enumerate(page['sections']):
                        placeholder = f'h1_{j}' if j < 10 else f'h2_{j-10}'
                        if self._replace_placeholder_with_pinyin(slide, placeholder, section, font_size=20):
                            pass
                        else:
                            # 如果没有对应占位符，尝试填充到已有位置
                            pass
                    
                    slides_to_keep.append(template_idx)
                    print(f"  目录已填充: {len(page['sections'])}个章节")
                else:
                    print("  警告: 没有目录模板，跳过")
            
            elif page['type'] == 'section':
                # 章节页
                if self.templates['section']:
                    template_idx = self.templates['section']['idx']
                    
                    if template_idx in slides_to_keep:
                        new_idx, slide = self._duplicate_slide(template_idx)
                        slides_to_keep.append(new_idx)
                    else:
                        slide = self.presentation.slides[template_idx]
                        slides_to_keep.append(template_idx)
                    
                    self._replace_placeholder_with_pinyin(slide, 'h1_0', page['title'], font_size=32)
                    self._clear_unused_placeholders(slide, ['h1_0'])
                    print(f"  章节页: {page['title']}")
            
            elif page['type'] == 'content':
                # 正文页
                text_count = len(page['contents'])
                template = self._find_best_content_template(text_count)
                
                if template:
                    template_idx = template['idx']
                    
                    if template_idx in slides_to_keep:
                        new_idx, slide = self._duplicate_slide(template_idx)
                        slides_to_keep.append(new_idx)
                    else:
                        slide = self.presentation.slides[template_idx]
                        slides_to_keep.append(template_idx)
                    
                    # 填充内容
                    used = ['h2_0']
                    self._replace_placeholder_with_pinyin(slide, 'h2_0', page['title'], font_size=28)
                    
                    for j, text in enumerate(page['contents']):
                        placeholder = f'h3_{j}'
                        self._replace_placeholder_with_pinyin(slide, placeholder, text, font_size=20)
                        used.append(placeholder)
                    
                    self._clear_unused_placeholders(slide, used)
                    print(f"  正文页: {page['title']} ({text_count}个文本)")
        
        # 添加结束页到保留列表
        if self.templates['end']:
            end_idx = self.templates['end']['idx']
            if end_idx not in slides_to_keep:
                slides_to_keep.append(end_idx)
        
        # 删除未使用的幻灯片
        all_indices = set(range(len(self.presentation.slides)))
        indices_to_remove = sorted(all_indices - set(slides_to_keep), reverse=True)
        
        print(f"\n删除 {len(indices_to_remove)} 个未使用的模板页...")
        for idx in indices_to_remove:
            rId = self.presentation.slides._sldIdLst[idx].rId
            self.presentation.part.drop_rel(rId)
            del self.presentation.slides._sldIdLst[idx]
        
        # 保存
        output_path = 'output.pptx'
        self.presentation.save(output_path)
        
        print(f"\n=== PPT生成完成 ===")
        print(f"输出文件: {output_path}")
        print(f"总页数: {len(self.presentation.slides)}")
        print(f"包含: 封面 + 目录 + 章节页 + 正文页 + 结束页")
        
        return output_path


if __name__ == '__main__':
    import sys
    from md_parser import MDParser
    
    if len(sys.argv) > 2:
        parser = MDParser(sys.argv[1])
        md_content = parser.parse()
        generator = PPTGenerator(sys.argv[2], md_content)
        generator.generate()
    else:
        print("用法: python ppt_generator.py <md文件> <模板文件>")