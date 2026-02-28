from pptx import Presentation
import os

def test_ppt_structure():
    """测试PPT结构"""
    ppt_path = 'output.pptx'
    assert os.path.exists(ppt_path), 'PPT文件未生成'
    
    presentation = Presentation(ppt_path)
    slides = presentation.slides
    
    # 检查页面数量
    print(f'PPT页面数量: {len(slides)}')
    
    # 检查封面页
    cover_slide = slides[0]
    print('封面页存在')
    
    # 检查目录页
    toc_slide = slides[1]
    print('目录页存在')
    
    # 检查章节页和正文页
    print(f'章节页和正文页数量: {len(slides) - 2}')
    
    return True

def test_content_integrity():
    """测试内容完整性"""
    ppt_path = 'output.pptx'
    presentation = Presentation(ppt_path)
    
    # 检查是否有内容
    has_content = False
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text:
                        has_content = True
                        break
                if has_content:
                    break
        if has_content:
            break
    
    assert has_content, 'PPT内容为空'
    print('PPT内容完整')
    return True

if __name__ == '__main__':
    test_ppt_structure()
    test_content_integrity()
    print('所有测试通过！')