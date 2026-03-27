from pptx import Presentation

def check_template():
    """检查模板文件的页面和占位符"""
    template_path = '模板.pptx'
    presentation = Presentation(template_path)
    
    print(f'模板页面数量: {len(presentation.slides)}')
    
    for i, slide in enumerate(presentation.slides):
        print(f'\n页面 {i+1}:')
        for j, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                print(f'  形状 {j+1} 文本:')
                for paragraph in shape.text_frame.paragraphs:
                    print(f'    {paragraph.text}')

if __name__ == '__main__':
    check_template()