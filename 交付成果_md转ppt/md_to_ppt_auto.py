import argparse
import copy
import os
import re
import tempfile
import urllib.request
from dataclasses import dataclass, field
from typing import List, Optional, Tuple

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


URL_RE = re.compile(r"https?://\S+")
MD_IMAGE_RE = re.compile(r"!\[[^\]]*\]\(([^)]+)\)")
PLACEHOLDER_RE = re.compile(r"\{\{([^}]+)\}\}")


@dataclass
class ContentBlock:
    section_title: str
    block_title: str
    texts: List[str] = field(default_factory=list)
    bullets: List[str] = field(default_factory=list)
    code_blocks: List[str] = field(default_factory=list)
    images: List[str] = field(default_factory=list)
    audios: List[str] = field(default_factory=list)
    videos: List[str] = field(default_factory=list)
    links: List[str] = field(default_factory=list)

    def all_text_lines(self) -> List[str]:
        lines: List[str] = []
        lines.extend(self.texts)
        lines.extend([f"• {item}" for item in self.bullets])
        lines.extend(self.code_blocks)
        lines.extend(self.audios)
        lines.extend(self.videos)
        lines.extend(self.links)
        return lines


@dataclass
class Section:
    title: str
    blocks: List[ContentBlock] = field(default_factory=list)


@dataclass
class ParsedDocument:
    title: str
    sections: List[Section]


@dataclass
class SlideProfile:
    template_index: int
    role: str
    text_slots: List[str]
    media_slots: List[str]
    image_shapes: int
    raw_placeholders: List[str]


class MarkdownParser:
    def __init__(self, md_path: str):
        self.md_path = md_path

    def parse(self) -> ParsedDocument:
        with open(self.md_path, "r", encoding="utf-8") as f:
            lines = f.read().splitlines()

        title = "未命名课程"
        sections: List[Section] = []
        current_section: Optional[Section] = None
        current_block: Optional[ContentBlock] = None

        in_code = False
        code_buffer: List[str] = []

        for raw in lines:
            line = raw.strip()

            if line.startswith("```"):
                if not in_code:
                    in_code = True
                    code_buffer = []
                else:
                    in_code = False
                    if current_block is not None and code_buffer:
                        current_block.code_blocks.append("\n".join(code_buffer))
                    code_buffer = []
                continue

            if in_code:
                code_buffer.append(raw.rstrip("\n"))
                continue

            if line.startswith("# "):
                title = line[2:].strip()
                continue

            if line.startswith("## "):
                current_section = Section(title=line[3:].strip())
                sections.append(current_section)
                current_block = None
                continue

            if line.startswith("### "):
                if current_section is None:
                    current_section = Section(title="默认章节")
                    sections.append(current_section)
                current_block = ContentBlock(
                    section_title=current_section.title,
                    block_title=line[4:].strip() or current_section.title,
                )
                current_section.blocks.append(current_block)
                continue

            if line.startswith("#### "):
                if current_block is None:
                    if current_section is None:
                        current_section = Section(title="默认章节")
                        sections.append(current_section)
                    current_block = ContentBlock(
                        section_title=current_section.title,
                        block_title=current_section.title,
                    )
                    current_section.blocks.append(current_block)
                text = line[5:].strip()
                if text:
                    current_block.texts.append(text)
                continue

            if not line:
                continue

            if current_section is None:
                current_section = Section(title="默认章节")
                sections.append(current_section)

            if current_block is None:
                current_block = ContentBlock(
                    section_title=current_section.title,
                    block_title=current_section.title,
                )
                current_section.blocks.append(current_block)

            image_match = MD_IMAGE_RE.search(line)
            if image_match:
                url = image_match.group(1).strip()
                self._dispatch_url(current_block, url)
                continue

            if line.startswith("- ") or line.startswith("* "):
                current_block.bullets.append(line[2:].strip())
                continue

            ordered_match = re.match(r"^\d+[.)]\s+(.*)$", line)
            if ordered_match:
                current_block.bullets.append(ordered_match.group(1).strip())
                continue

            urls = URL_RE.findall(line)
            if urls and len(line.replace(urls[0], "").strip()) == 0:
                for url in urls:
                    self._dispatch_url(current_block, url)
                continue

            current_block.texts.append(line)

        return ParsedDocument(title=title, sections=sections)

    def _dispatch_url(self, block: ContentBlock, url: str) -> None:
        media_type = self._detect_media_type(url)
        if media_type == "image":
            block.images.append(url)
        elif media_type == "audio":
            block.audios.append(url)
        elif media_type == "video":
            block.videos.append(url)
        else:
            block.links.append(url)

    @staticmethod
    def _detect_media_type(url: str) -> str:
        lower = url.lower().split("?")[0]
        if lower.endswith((".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp")):
            return "image"
        if lower.endswith((".mp3", ".wav", ".m4a", ".aac", ".ogg")):
            return "audio"
        if lower.endswith((".mp4", ".mov", ".avi", ".mkv", ".webm")):
            return "video"
        return "link"


class TemplateAnalyzer:
    def __init__(self, template_path: str):
        self.prs = Presentation(template_path)
        self.profiles = self._build_profiles()

    def _build_profiles(self) -> List[SlideProfile]:
        profiles: List[SlideProfile] = []
        for idx, slide in enumerate(self.prs.slides):
            placeholders = self._extract_placeholders(slide)
            texts = [p for p in placeholders if p.startswith("h3_")]
            medias = [p for p in placeholders if p.startswith(("h4_", "h5_", "h6_"))]
            image_shapes = sum(1 for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE)

            role = "other"
            if "h0_0" in placeholders:
                role = "cover"
            elif "h1_0" in placeholders and "h2_0" not in placeholders:
                role = "section"
            elif "h2_0" in placeholders:
                role = "content"
            else:
                all_text = self._slide_text(slide)
                if "目录" in all_text or "mù lù" in all_text:
                    role = "toc"
                if "xiè xie" in all_text or "谢" in all_text:
                    role = "end"

            profiles.append(
                SlideProfile(
                    template_index=idx,
                    role=role,
                    text_slots=sorted(texts),
                    media_slots=sorted(medias),
                    image_shapes=image_shapes,
                    raw_placeholders=sorted(placeholders),
                )
            )
        return profiles

    @staticmethod
    def _extract_placeholders(slide) -> List[str]:
        tokens: List[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    tokens.extend(PLACEHOLDER_RE.findall(paragraph.text))
        return list(dict.fromkeys(tokens))

    @staticmethod
    def _slide_text(slide) -> str:
        chunks: List[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                chunks.append(shape.text)
        return "\n".join(chunks)


class MarkdownToPPT:
    def __init__(self, md_path: str, template_path: str, output_path: str):
        self.md_path = md_path
        self.template_path = template_path
        self.output_path = output_path

        self.parsed = MarkdownParser(md_path).parse()
        self.analyzer = TemplateAnalyzer(template_path)
        self.template_prs = self.analyzer.prs
        self.output_prs = Presentation(template_path)
        self.temp_files: List[str] = []

        self.slide_plan: List[Tuple[str, Optional[ContentBlock], Optional[str], int]] = []

    def build(self) -> str:
        self._plan_slides()
        self._render()
        self.output_prs.save(self.output_path)
        self._cleanup_temp_files()
        return self.output_path

    def _plan_slides(self) -> None:
        cover = [p for p in self.analyzer.profiles if p.role == "cover"]
        toc = [p for p in self.analyzer.profiles if p.role == "toc"]
        section = [p for p in self.analyzer.profiles if p.role == "section"]
        content = [p for p in self.analyzer.profiles if p.role == "content"]
        end = [p for p in self.analyzer.profiles if p.role == "end"]

        cover_idx = cover[0].template_index if cover else 0
        toc_idx = toc[0].template_index if toc else cover_idx
        section_idx = section[0].template_index if section else cover_idx
        end_idx = end[0].template_index if end else -1

        self.slide_plan.append(("cover", None, None, cover_idx))

        intro_blocks: List[ContentBlock] = []
        normal_sections = self.parsed.sections
        if self.parsed.sections:
            first = self.parsed.sections[0]
            if first.blocks and ("今天学什么" in first.title or "导入" in first.title or len(first.blocks) >= 2):
                intro_blocks = first.blocks
                normal_sections = self.parsed.sections[1:]

        for block in intro_blocks:
            self.slide_plan.append(("content", block, None, self._pick_content_slide(content, block)))

        self.slide_plan.append(("toc", None, None, toc_idx))

        for sec in normal_sections:
            self.slide_plan.append(("section", None, sec.title, section_idx))
            for block in sec.blocks:
                self.slide_plan.append(("content", block, None, self._pick_content_slide(content, block)))

        if end_idx >= 0:
            self.slide_plan.append(("end", None, None, end_idx))

    @staticmethod
    def _pick_content_slide(candidates: List[SlideProfile], block: ContentBlock) -> int:
        if not candidates:
            return 0

        need_text = max(1, len(block.all_text_lines()))
        need_images = len(block.images)
        need_media = len(block.audios) + len(block.videos) + len(block.links)

        best = candidates[0]
        best_score = -10**9

        for c in candidates:
            text_slots = len(c.text_slots) if c.text_slots else 1
            score = 0

            if text_slots >= need_text:
                score += 100 - (text_slots - need_text) * 5
            else:
                score -= (need_text - text_slots) * 30

            if need_images >= 3 and c.image_shapes >= 3:
                score += 20
            if need_images == 0 and c.image_shapes >= 4:
                score -= 12
            if need_media > 0 and text_slots >= 2:
                score += 8
            if len(block.code_blocks) > 0 and text_slots >= 2:
                score += 10

            if score > best_score:
                best_score = score
                best = c

        return best.template_index

    def _render(self) -> None:
        while len(self.output_prs.slides) > 0:
            self._remove_slide(self.output_prs, 0)

        for role, block, sec_title, template_idx in self.slide_plan:
            slide = self._duplicate_template_slide(template_idx)

            if role == "cover":
                self._fill_placeholder(slide, "h0_0", self.parsed.title)
            elif role == "toc":
                for i, sec in enumerate(self.parsed.sections):
                    self._fill_placeholder(slide, f"h1_{i}", sec.title)
            elif role == "section" and sec_title is not None:
                self._fill_placeholder(slide, "h1_0", sec_title)
            elif role == "content" and block is not None:
                self._fill_content_slide(slide, block)

            self._clear_unused_placeholders(slide)

    def _duplicate_template_slide(self, template_idx: int):
        source = self.template_prs.slides[template_idx]
        dest = self.output_prs.slides.add_slide(self.output_prs.slide_layouts[6])

        for shape in source.shapes:
            new_el = copy.deepcopy(shape.element)
            dest.shapes._spTree.insert_element_before(new_el, "p:extLst")

        for rel in source.part.rels.values():
            if "notesSlide" in rel.reltype:
                continue
            dest.part.rels._add_relationship(rel.reltype, rel._target, rel.rId)

        return dest

    @staticmethod
    def _remove_slide(prs: Presentation, idx: int) -> None:
        slide_id = prs.slides._sldIdLst[idx]
        prs.slides._sldIdLst.remove(slide_id)

    def _fill_content_slide(self, slide, block: ContentBlock) -> None:
        self._fill_placeholder(slide, "h2_0", block.block_title)

        lines = block.all_text_lines()
        for i, line in enumerate(lines):
            self._fill_placeholder(slide, f"h3_{i}", line)

        for i, url in enumerate(block.images):
            if not self._fill_placeholder(slide, f"h4_{i}", url):
                self._fill_image_to_picture_slot(slide, i, url)

        for i, url in enumerate(block.audios):
            if not self._fill_placeholder(slide, f"h5_{i}", url):
                self._append_to_text_slot(slide, url)

        for i, url in enumerate(block.videos):
            if not self._fill_placeholder(slide, f"h6_{i}", url):
                self._append_to_text_slot(slide, url)

        for url in block.links:
            if not self._fill_placeholder(slide, "h5_0", url):
                self._append_to_text_slot(slide, url)

    def _fill_image_to_picture_slot(self, slide, image_idx: int, url: str) -> None:
        picture_shapes = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        if image_idx >= len(picture_shapes):
            self._append_to_text_slot(slide, url)
            return

        image_path = self._download_url(url)
        if not image_path:
            self._append_to_text_slot(slide, url)
            return

        pic = picture_shapes[image_idx]
        left, top, width, height = pic.left, pic.top, pic.width, pic.height
        pic.element.getparent().remove(pic.element)
        slide.shapes.add_picture(image_path, left, top, width, height)

    def _download_url(self, url: str) -> Optional[str]:
        if not url.startswith("http"):
            return None
        try:
            suffix = ".jpg"
            lower = url.lower().split("?")[0]
            if "." in lower:
                ext = "." + lower.rsplit(".", 1)[1]
                if len(ext) <= 6:
                    suffix = ext
            fd, path = tempfile.mkstemp(suffix=suffix)
            os.close(fd)
            urllib.request.urlretrieve(url, path)
            self.temp_files.append(path)
            return path
        except Exception:
            return None

    def _cleanup_temp_files(self) -> None:
        for path in self.temp_files:
            try:
                if os.path.exists(path):
                    os.remove(path)
            except Exception:
                pass

    @staticmethod
    def _find_runs_with_placeholder(slide, token: str):
        marker = "{{" + token + "}}"
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if marker in run.text:
                            yield run

    def _fill_placeholder(self, slide, token: str, value: str) -> bool:
        replaced = False
        marker = "{{" + token + "}}"
        for run in self._find_runs_with_placeholder(slide, token):
            run.text = run.text.replace(marker, value)
            replaced = True
        return replaced

    def _append_to_text_slot(self, slide, value: str) -> None:
        for key in ["h3_0", "h3_1", "h3_2", "h3_3"]:
            marker = "{{" + key + "}}"
            for run in self._find_runs_with_placeholder(slide, key):
                run.text = run.text.replace(marker, value)
                return

    @staticmethod
    def _clear_unused_placeholders(slide) -> None:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = PLACEHOLDER_RE.sub("", run.text)


def main() -> None:
    parser = argparse.ArgumentParser(description="Markdown 自动匹配模板生成 PPT")
    parser.add_argument("--md", required=True, help="Markdown 文件路径")
    parser.add_argument("--template", required=True, help="模板 PPT 路径")
    parser.add_argument("--out", required=True, help="输出 PPT 路径")
    args = parser.parse_args()

    generator = MarkdownToPPT(args.md, args.template, args.out)
    out = generator.build()
    print(f"生成完成: {out}")


if __name__ == "__main__":
    main()
